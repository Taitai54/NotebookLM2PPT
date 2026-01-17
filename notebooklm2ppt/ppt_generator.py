"""
PowerPoint Generator for NotebookLM slide reconstruction.

Generates fully editable PPTX presentations from separated slide layers
(background images, text blocks, and graphics).
"""

from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
import cv2
import os
from typing import List, Dict, Tuple, Optional

# Import configuration
from .config import (
    SLIDE_WIDTH_INCHES,
    SLIDE_HEIGHT_INCHES,
    EMUS_PER_INCH,
)

# Import coordinate utilities
from .utils.coordinates import (
    pixels_to_pptx_coordinates,
    validate_bbox_in_bounds,
)


class PPTCreator:
    """
    Generate fully editable PPTX from separated layers.
    
    Creates slides with:
    - Original images as slide backgrounds
    - Extracted text as editable text boxes
    - Proper positioning and formatting preserved
    """
    
    def __init__(
        self,
        width_inches: float = SLIDE_WIDTH_INCHES,
        height_inches: float = SLIDE_HEIGHT_INCHES
    ):
        """
        Initialize the PowerPoint creator.
        
        Args:
            width_inches: Slide width in inches (default: 16 for 16:9)
            height_inches: Slide height in inches (default: 9 for 16:9)
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(width_inches)
        self.prs.slide_height = Inches(height_inches)
        self.slides_created = 0
        
        # Store dimensions in EMUs for coordinate conversion
        self.slide_width_emu = int(width_inches * EMUS_PER_INCH)
        self.slide_height_emu = int(height_inches * EMUS_PER_INCH)

    def _add_background_image(
        self,
        slide,
        image_path: str
    ) -> None:
        """
        Add image as slide background at full bleed.
        """
        slide.shapes.add_picture(
            image_path, 
            0, 0, 
            self.prs.slide_width, 
            self.prs.slide_height
        )

    def _calculate_font_size(
        self,
        font_size_px: int,
        img_height: int
    ) -> float:
        """
        Calculate font size in points from pixel height.
        
        Uses a heuristic scaling based on image dimensions.
        """
        # Heuristic: scale based on image height to slide height ratio
        # Then apply adjustment factor (0.95) to prevent overflow
        pt_size = (font_size_px / img_height) * 7.5 * 72 * 0.95
        
        # Clamp to reasonable range
        if pt_size < 8:
            pt_size = 8
        elif pt_size > 144:
            pt_size = 144
            
        return pt_size

    def _add_text_box(
        self,
        slide,
        text: str,
        bbox: Tuple[int, int, int, int],
        font_size_px: int,
        scale_x: float,
        scale_y: float,
        role: str = "body"
    ) -> None:
        """
        Add an editable text box to the slide.
        
        Args:
            slide: The slide to add to
            text: Text content
            bbox: (x, y, width, height) in pixels
            font_size_px: Font size in pixels
            scale_x: X scale factor (slide width / image width)
            scale_y: Y scale factor (slide height / image height)
            role: Text role (title, body, etc.)
        """
        x, y, w, h = bbox
        
        # Convert to PPT coordinates using scale factors
        ppt_x = int(x * scale_x)
        ppt_y = int(y * scale_y)
        ppt_w = int(w * scale_x)
        ppt_h = int(h * scale_y)
        
        # Create text box
        textbox = slide.shapes.add_textbox(ppt_x, ppt_y, ppt_w, ppt_h)
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        
        # Set text
        p = tf.paragraphs[0]
        p.text = text
        
        # Calculate and set font size
        try:
            # Get image height from scale factor
            img_h = int(self.slide_height_emu / scale_y) if scale_y > 0 else 1080
            pt_size = self._calculate_font_size(font_size_px, img_h)
            p.font.size = Pt(pt_size)
        except Exception:
            # Fallback to default size
            pass

    def _add_image_object(
        self,
        slide,
        img_obj: Dict,
        scale_x: float,
        scale_y: float
    ) -> None:
        """
        Add an extracted image object (diagram, icon) to the slide.
        """
        img_path = img_obj.get('path', '')
        if not img_path or not os.path.exists(img_path):
            return
            
        x, y, w, h = img_obj.get('box', [0, 0, 100, 100])
        
        ppt_x = int(x * scale_x)
        ppt_y = int(y * scale_y)
        ppt_w = int(w * scale_x)
        ppt_h = int(h * scale_y)
        
        try:
            slide.shapes.add_picture(img_path, ppt_x, ppt_y, ppt_w, ppt_h)
        except Exception as e:
            print(f"Warning: Could not add image object {img_path}: {e}")

    def add_slide(
        self,
        background_image_path: str,
        text_blocks: List[Dict],
        image_objects: List[Dict],
        original_image_size: Tuple[int, int]
    ) -> None:
        """
        Add a reconstructed slide to the presentation.
        
        Args:
            background_image_path: Path to the clean background image
            text_blocks: List of text paragraphs with box, text, font_size
            image_objects: List of extracted image objects (diagrams)
            original_image_size: Tuple (width, height) of the source image
        """
        blank_slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(blank_slide_layout)
        
        # 1. Add Background
        self._add_background_image(slide, background_image_path)
        
        # Calculate scale factors
        img_w, img_h = original_image_size
        scale_x = self.prs.slide_width / img_w
        scale_y = self.prs.slide_height / img_h
        
        # 2. Add Extracted Image Objects (Diagrams/Photos)
        for img_obj in image_objects:
            self._add_image_object(slide, img_obj, scale_x, scale_y)
        
        # 3. Add Text Boxes
        for block in text_blocks:
            text = block.get('text', '')
            bbox = block.get('box', [0, 0, 100, 20])
            font_size_px = block.get('font_size', 20)
            role = block.get('role', 'body')
            
            self._add_text_box(
                slide, text, tuple(bbox), font_size_px,
                scale_x, scale_y, role
            )
        
        self.slides_created += 1

    def save(self, output_path: str) -> None:
        """
        Save presentation to PPTX file.
        """
        self.prs.save(output_path)
        print(f"✓ Saved {self.slides_created} slides to {output_path}")


class PowerPointGenerator(PPTCreator):
    """
    Alias for PPTCreator to match spec naming.
    
    Provides the same functionality with the name used in the spec.
    """
    pass
