from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
import cv2
import os

class PPTCreator:
    def __init__(self):
        self.prs = Presentation()
        # Set slide size to 16:9 
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)

    def add_slide(self, background_image_path, text_blocks, image_objects, original_image_size):
        """
        Add a reconstructed slide to the presentation.
        
        Args:
            background_image_path: Path to the clean background image
            text_blocks: List of text paragraphs
            image_objects: List of extracted image objects (diagrams)
            original_image_size: Tuple (width, height) of the image for scaling
        """
        blank_slide_layout = self.prs.slide_layouts[6] 
        slide = self.prs.slides.add_slide(blank_slide_layout)
        
        # 1. Add Background
        slide.shapes.add_picture(
            background_image_path, 
            0, 0, 
            self.prs.slide_width, 
            self.prs.slide_height
        )
        
        img_w, img_h = original_image_size
        scale_x = self.prs.slide_width / img_w
        scale_y = self.prs.slide_height / img_h
        
        # 2. Add Extracted Image Objects (Diagrams/Photos)
        for img_obj in image_objects:
            img_path = img_obj['path']
            x, y, w, h = img_obj['box']
            
            ppt_x = int(x * scale_x)
            ppt_y = int(y * scale_y)
            ppt_w = int(w * scale_x)
            ppt_h = int(h * scale_y)
            
            try:
                slide.shapes.add_picture(img_path, ppt_x, ppt_y, ppt_w, ppt_h)
            except Exception as e:
                print(f"Warning: Could not add image object {img_path}: {e}")
        
        # 3. Add Text Boxes
        
        for block in text_blocks:
            text = block['text']
            x, y, w, h = block['box']
            font_size_px = block.get('font_size', 20)
            
            # Convert to PPT coordinates
            ppt_x = int(x * scale_x)
            ppt_y = int(y * scale_y)
            # Ensure width is at least something substantial to prevent tiny wrapping
            ppt_w = max(int(w * scale_x), int(1 * 914400)) # Min 1 inch width? No, that's too much.
            ppt_w = int(w * scale_x)
            ppt_h = int(h * scale_y)
            
            # Create text box
            textbox = slide.shapes.add_textbox(ppt_x, ppt_y, ppt_w, ppt_h)
            tf = textbox.text_frame
            tf.word_wrap = True
            # CRITICAL FIX: Allow shape to expand if text finds it needs more room
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            
            p = tf.paragraphs[0]
            p.text = text
            
            try:
                # Heuristic: 0.75 ratio often works for converting pixel height to point size
                # Slightly reducing it to 0.7 to prevent "too large" text overflowing
                pt_size = (font_size_px / img_h) * 7.5 * 72 * 0.95 
                if pt_size < 8: pt_size = 8
                
                p.font.size = Pt(pt_size)
            except:
                pass

    def save(self, path):
        self.prs.save(path)
