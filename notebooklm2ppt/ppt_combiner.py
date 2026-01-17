# Merge PPT files using Spire.Presentation, preserving original design
# Install: pip install spire.presentation

import os
from spire.presentation import *
from spire.presentation.common import *
from pptx import Presentation as PptxPresentation

def combine_ppt_files_with_spire(source_folder, output_file):
    """
    Merge PPT files using Spire.Presentation, keeping only the first page of each PPT and preserving original design
    
    Args:
        source_folder: Path to the folder containing source PPT files
        output_file: Output path for the merged PPT file
    """
    # Get all pptx files and sort alphabetically
    ppt_files = sorted([f for f in os.listdir(source_folder) if f.endswith('.pptx')])
    
    if not ppt_files:
        print("No PPT files found")
        return
    
    print(f"Found {len(ppt_files)} PPT files:")
    for idx, file in enumerate(ppt_files, 1):
        print(f"  {idx}. {file}")
    
    # Create main presentation object, using the first PPT as the base
    first_ppt_path = os.path.join(source_folder, ppt_files[0])
    main_pres = Presentation()
    main_pres.LoadFromFile(first_ppt_path)
    
    # Delete extra pages from the first PPT, keeping only the first page
    while main_pres.Slides.Count > 1:
        main_pres.Slides.RemoveAt(1)
    
    print(f"  Added: {ppt_files[0]} (page 1)")
    
    # Iterate through remaining PPT files
    for ppt_file in ppt_files[1:]:
        file_path = os.path.join(source_folder, ppt_file)
        
        # Load current PPT file
        temp_pres = Presentation()
        temp_pres.LoadFromFile(file_path)
        
        if temp_pres.Slides.Count > 0:
            # Use AppendBySlide method to append the first page, preserving original design
            main_pres.Slides.AppendBySlide(temp_pres.Slides[0])
            print(f"  Added: {ppt_file} (page 1)")
        else:
            print(f"  Skipped: {ppt_file} (no slides)")
        
        # Release temporary presentation resources
        temp_pres.Dispose()
    
    # Save merged PPT
    main_pres.SaveToFile(output_file, FileFormat.Pptx2016)
    print(f"\nMerge complete! Output file: {output_file}")
    print(f"Total merged slides: {main_pres.Slides.Count}")
    
    # Release resources
    main_pres.Dispose()


def combine_ppt_files_with_master(source_folder, output_file):
    """
    Merge PPT files using Spire.Presentation, using a unified master design
    
    Args:
        source_folder: Path to the folder containing source PPT files
        output_file: Output path for the merged PPT file
    """
    # Get all pptx files and sort alphabetically
    ppt_files = sorted([f for f in os.listdir(source_folder) if f.endswith('.pptx')])
    
    if not ppt_files:
        print("No PPT files found")
        return
    
    print(f"Found {len(ppt_files)} PPT files:")
    for idx, file in enumerate(ppt_files, 1):
        print(f"  {idx}. {file}")
    
    # Create main presentation object, using the first PPT as the base
    first_ppt_path = os.path.join(source_folder, ppt_files[0])
    main_pres = Presentation()
    main_pres.LoadFromFile(first_ppt_path)
    
    # Delete extra pages from the first PPT, keeping only the first page
    while main_pres.Slides.Count > 1:
        main_pres.Slides.RemoveAt(1)
    
    print(f"  Added: {ppt_files[0]} (page 1)")
    
    # Get the master from the first presentation
    master = main_pres.Masters[0]
    
    # Iterate through remaining PPT files
    for ppt_file in ppt_files[1:]:
        file_path = os.path.join(source_folder, ppt_file)
        
        # Load current PPT file
        temp_pres = Presentation()
        temp_pres.LoadFromFile(file_path)
        
        if temp_pres.Slides.Count > 0:
            # Use AppendByMaster method to append the first page, applying unified master
            main_pres.Slides.AppendByMaster(temp_pres.Slides[0], master)
            print(f"  Added: {ppt_file} (page 1, using unified master)")
        else:
            print(f"  Skipped: {ppt_file} (no slides)")
        
        # Release temporary presentation resources
        temp_pres.Dispose()
    
    # Save merged PPT
    main_pres.SaveToFile(output_file, FileFormat.Pptx2016)
    print(f"\nMerge complete! Output file: {output_file}")
    print(f"Total merged slides: {main_pres.Slides.Count}")
    
    # Release resources
    main_pres.Dispose()

def combine_ppt(source_folder, out_ppt_file):
    # Ensure string paths, as we use .replace later
    source_folder = str(source_folder)
    out_ppt_file = str(out_ppt_file)
    
    # Method 1: Preserve original design (recommended)
    output_file1 = out_ppt_file.replace(".pptx", "_combined_original_design.pptx")
    print("=" * 60)
    print("Method 1: Merge PPT and preserve original design")
    print("=" * 60)
    combine_ppt_files_with_spire(source_folder, output_file1)


    ppt = PptxPresentation(output_file1)
    for slide in ppt.slides:
        for shape in list(slide.shapes):
            # If shape name is "New shape", delete it
            if shape.name == "New shape":
                sp = slide.shapes._spTree.remove(shape._element)

    ppt.save(out_ppt_file)
    print(f"Deleted temporary file: {output_file1}")
    os.remove(output_file1)
